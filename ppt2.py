from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt

class PPTGenerator:
    def __init__(self, contents) -> None:
        self.contents = contents

    def addTitleSlide(self, prs: Presentation) -> None:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = self.contents['heading']
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Red color
        prs.save(self.contents["Presentation Name"])

    def addTableOfContents(self, prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Contents"
        title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Red color
        bullet_points = slide.shapes.placeholders[1]

        for subheading_dict in self.contents["subheadings"]:
            subheading = subheading_dict["subheading"]
            text_frame = bullet_points.text_frame
            p = text_frame.add_paragraph()
            p.text = subheading
            p.level = 0
            p.runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Blue color

        prs.save(self.contents["Presentation Name"])

    def addTopicExplanation(self, prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]
        slides = []

        for subheading_dict in self.contents["subheadings"]:
            subheading = subheading_dict["subheading"]
            content = subheading_dict["content"]

            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = subheading
            title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(88, 24, 69)  # Red color
            bullet_points = slide.shapes.placeholders[1]

            current_slide = slide
            current_frame = bullet_points.text_frame

            for explanation in content:
                p = current_frame.add_paragraph()
                p.text = explanation
                p.level = 0
                p.runs[0].font.color.rgb = RGBColor(144, 12, 63)  # Blue color

                if self._is_overflowing(current_frame):
                    slides.append(current_slide)
                    slide = prs.slides.add_slide(slide_layout)
                    title_shape = slide.shapes.title
                    title_shape.text = subheading
                    title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(88, 24, 69)  # Red color
                    bullet_points = slide.shapes.placeholders[1]
                    current_slide = slide
                    current_frame = bullet_points.text_frame
                    p = current_frame.add_paragraph()
                    p.text = explanation
                    p.level = 0
                    p.runs[0].font.color.rgb = RGBColor(144, 12, 63)  # Blue color

            slides.append(current_slide)

        # Remove empty slides
        for slide in slides:
            if not any(shape.has_text_frame for shape in slide.shapes):
                prs.slides._sldIdLst.remove(slide._element)

        prs.save(self.contents["Presentation Name"])



    def _is_overflowing(self, text_frame):
        text_content = text_frame.text
        text_height = Pt(len(text_content))  # Measure the height of the text content
        max_height = Inches(5)  # Adjust as needed
        return text_height > max_height



if __name__ == "__main__":
    import json

    # Load contents from the JSON file
    with open("static/1.json", "r") as file:
        content = json.load(file)

    newContent = {
        "Presentation Name": "static/Sample.pptx",
        "heading": content["heading"],
        "subheadings": content["subheadings"]
    }

    MyPptGenerator = PPTGenerator(contents=newContent)
    prs = Presentation()
    MyPptGenerator.addTitleSlide(prs)
    MyPptGenerator.addTableOfContents(prs)
    MyPptGenerator.addTopicExplanation(prs)