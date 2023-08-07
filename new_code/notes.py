from collections.abc import Container
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

def add_notes_to_slide(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes

# Replace these variables with your actual values
slide_index = 0   # The index of the slide you want to add notes to (0-based index)
notes = "These are the speaker notes for the first slide."

# Load the presentation
presentation = Presentation("output.pptx")

# Get the slide you want to add notes to
slide = presentation.slides[slide_index]

# Call the function to add the notes to the slide
add_notes_to_slide(slide, notes)

# Save the updated presentation with notes
presentation.save("output_with_notes.pptx")
