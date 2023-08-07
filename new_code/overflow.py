from collections.abc import Container
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from pptx import Presentation

def split_text_between_two_slides(slide, max_lines_per_slide=5):
    if slide is None:
        return

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            paragraphs = list(text_frame.paragraphs)
            total_paragraphs = len(paragraphs)

            if total_paragraphs > max_lines_per_slide:
                # Create a new slide to continue the text
                new_slide_layout = slide.slide_layout
                presentation = slide.presentation
                new_slide = presentation.slides.add_slide(new_slide_layout)

                # Copy the remaining paragraphs to the new slide
                for para in paragraphs[max_lines_per_slide:]:
                    new_text_frame = new_slide.shapes.add_textbox(0, 0, presentation.slide_width, presentation.slide_height).text_frame
                    new_text_frame.text = para.text

                # Remove the remaining paragraphs from the original slide
                for _ in range(total_paragraphs - max_lines_per_slide):
                    text_frame.paragraphs[-1].element.getparent().remove(text_frame.paragraphs[-1].element)

                # Adjust vertical alignment of the original slide text box
                text_frame.vertical_anchor = text_frame.vertical_anchor - 1

                break

# Load the PowerPoint presentation
presentation_path = 'Sample.pptx'
presentation = Presentation(presentation_path)

# Specify the slide number you want to split the text
slide_number_to_split = 3

# Define the maximum number of lines per slide
max_lines_per_slide = 5

# Get the slide to split the text
slide_to_split = presentation.slides[slide_number_to_split - 1]

# Split text between two slides for the specified slide
split_text_between_two_slides(slide_to_split, max_lines_per_slide)

# Save the modified presentation
presentation.save('modified_presentation.pptx')


