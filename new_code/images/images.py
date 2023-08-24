import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

def add_image_from_url_to_placeholder(slide, placeholder_idx, url):
    response = requests.get(url)
    if response.status_code == 200:
        image_stream = BytesIO(response.content)
        # slide.shapes._spTree.remove(slide.placeholders[placeholder_idx]._element)
        
        left = Inches(1)   # Adjust the left, top, width, and height values as needed
        top = Inches(1)
        width = Inches(6)
        height = Inches(4)

        slide.shapes.add_picture(image_stream, left, top, width, height)
    else:
        print(f"Failed to fetch the image from URL: {url}")

# Replace these variables with your actual values
url = "https://i0.wp.com/i.ytimg.com/vi/bAJ5hfE6wgs/hqdefault.jpg?w=780&ssl=1"
placeholder_idx = 1   # Adjust this value based on the index of the image placeholder

# Create a presentation and a slide to add the image
presentation = Presentation()
slide_layout = presentation.slide_layouts[8]  # Use an appropriate slide layout here
for i, shape in enumerate(slide_layout.placeholders):
    print(f"Placeholder index {i}: {shape.name}")

slide = presentation.slides.add_slide(presentation.slide_layouts[4])  # Use an appropriate slide layout here

# Call the function to add the image from the URL to the designated image box
add_image_from_url_to_placeholder(slide, placeholder_idx, url)

# Save the presentation
presentation.save("output.pptx")
