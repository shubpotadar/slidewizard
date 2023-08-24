from moviepy.editor import VideoFileClip

video_path = 'giphy.mp4'
gif_output_path = 'output.gif'

video_clip = VideoFileClip(video_path)
video_clip.write_gif(gif_output_path)

from pptx import Presentation
from pptx.util import Inches

presentation = Presentation()

gif_path = 'output.gif'  # Path to the converted GIF

slide = presentation.slides.add_slide(presentation.slide_layouts[5])
left = Inches(1)
top = Inches(1)
width = Inches(8)
height = Inches(6)
pic = slide.shapes.add_picture(gif_path, left, top, width, height)

presentation.save('output.pptx')
