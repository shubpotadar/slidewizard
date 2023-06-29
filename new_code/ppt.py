from collections.abc import Container
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

class PPTGenerator:
    def __init__(self, contents) -> None:
        self.contents = contents

    def addTitleSlide(self, prs: Presentation) -> None:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = self.contents['heading']
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Red color
        prs.save(self.contents["Presentation Name"])

    def addTableOfContents(self, prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Contents"
        title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Red color
        bullet_points = slide.shapes.placeholders[1]

        for subheading_dict in self.contents["subheadings"]:
            subheading = subheading_dict["subheading"]
            text_frame = bullet_points.text_frame
            p = text_frame.add_paragraph()
            p.text = subheading
            p.level = 0
            p.runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Blue color

        prs.save(self.contents["Presentation Name"])

    def addTopicExplanation(self, prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]

        for subheading_dict in self.contents["subheadings"]:
            subheading = subheading_dict["subheading"]
            content = subheading_dict["content"]

            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = subheading
            title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor( 88, 24, 69 )  # Red color
            bullet_points = slide.shapes.placeholders[1]

            for explanation in content:
                text_frame = bullet_points.text_frame
                p = text_frame.add_paragraph()
                p.text = explanation
                p.level = 0
                p.runs[0].font.color.rgb = RGBColor(0, 0, 255)  # Blue color

        prs.save(self.contents["Presentation Name"])


if __name__ == "__main__":
    import json
    import os
    import pptx.util

    # Load contents from the JSON file
    with open("./1.json", "r") as file:
        content = json.load(file)
    cwd = os.getcwd()

    newContent = {
        "Presentation Name": cwd + "\\Sample.pptx",
        "heading": content["heading"],
        "subheadings": content["subheadings"]
    }

    MyPptGenerator = PPTGenerator(contents=newContent)
    prs = Presentation("title.pptx")
    for i in range(len(prs.slides)-1, -1, -1): 
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    MyPptGenerator.addTitleSlide(prs)
    MyPptGenerator.addTableOfContents(prs)
    MyPptGenerator.addTopicExplanation(prs)       